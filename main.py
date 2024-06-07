from fastapi import FastAPI, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from pathlib import Path

from fastapi.responses import FileResponse
import shutil

# импортируем модуль sumy
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer
from sumy.nlp.stemmers import Stemmer
from sumy.utils import get_stop_words

import os
import ffmpeg


from pydub import AudioSegment

import whisper_model

from pptx import Presentation
from pptx.util import Inches

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


UPLOAD_DIR = Path() / 'uploads'

# задаем язык и количество предложений в резюме
LANGUAGE = 'russian'
SENTENCES_COUNT = 30


app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=['*'],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.post("/uploadfile/")
async def create_upload_files(file_upload: UploadFile):
    data = await file_upload.read()
    save_to = UPLOAD_DIR / file_upload.filename
    with open(save_to, 'wb') as f:
        f.write(data)

    audio_path = await convert_video_to_audio(file_upload.filename)
    if audio_path:
        text_path = await convert_audio_to_text(audio_path)

        return FileResponse(UPLOAD_DIR / 'result.pptx', media_type='application/octet-stream', filename='result.pptx')
    
    


async def convert_video_to_audio(filename):
    video_path = os.path.join(UPLOAD_DIR, filename)
    audio_path = os.path.join(UPLOAD_DIR, os.path.splitext(filename)[0] + '.mp3')

    try:
        # Extract audio from the video file
        stream = ffmpeg.input(video_path)
        stream = ffmpeg.output(stream, audio_path)
        ffmpeg.run(stream)
        print(f"Audio file saved: {audio_path}")
        return audio_path
    except Exception as e:
        print(f"Error converting video to audio: {e}")
        return None
    


async def convert_audio_to_text(audio_path):
# Load the audio file
    audio = AudioSegment.from_file(audio_path, format="mp3")

    # Save the audio to a WAV file (required for speech_recognition)
    wav_path = os.path.splitext(audio_path)[0] + ".wav"
    audio.export(wav_path, format="wav")
    
    result = whisper_model.model.transcribe(wav_path) #результат работы транскрибирования из аудиофайла
    with open(UPLOAD_DIR / "tutor.txt", "w", encoding="utf-8") as f:
        f.write(result["text"]) #записываем в новый текстовый файл результат работы программы (только текст)

    parser = PlaintextParser.from_file(UPLOAD_DIR / "tutor.txt", Tokenizer(LANGUAGE))
    stemmer = Stemmer(LANGUAGE)
    summarizer = LsaSummarizer(stemmer)
    summarizer.stop_words = get_stop_words(LANGUAGE)
   
   # выводим резюме
    result_path = UPLOAD_DIR / "result.txt"
    with open(result_path, "w", encoding="utf-8") as f:
        for index, sentence in enumerate(summarizer(parser.document, SENTENCES_COUNT)):
            f.write(str(index+1) + ". " + str(sentence)+"\n")


    presentation = Presentation()
    with open(result_path, 'r', encoding="utf-8") as file:
        lines = file.readlines()
        for line in lines:
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
                    # Добавляем фон
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(204, 229, 255)  # Светло-голубой цвет фона
            
            # Добавляем текст в подзаголовок
            slide.shapes.title.text = ''  # Чистим заголовок
            slide.shapes.placeholders[1].text = line.strip()  # Записываем строку в подзаголовок
            
            # Форматируем текст в подзаголовке
            placeholder = slide.shapes.placeholders[1]
            placeholder.text_frame.paragraphs[0].font.size = Pt(18)
            placeholder.text_frame.paragraphs[0].font.bold = True
            placeholder.text_frame.paragraphs[0].font.name = 'Arial'
            placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Выравнивание по центру
            
            # Добавляем рамку вокруг подзаголовка
            placeholder.line.color.rgb = RGBColor(0, 0, 0)  # Черный цвет рамки
            placeholder.line.width = Pt(2)  # Толщина рамки 2 пункта

    pptx_filename = os.path.splitext(os.path.basename(result_path))[0] + ".pptx"
    save_to = UPLOAD_DIR / pptx_filename
    presentation.save(save_to)
    


    

